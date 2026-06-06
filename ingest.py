import os
import uuid

from pptx import Presentation
from openai import OpenAI
from pinecone import Pinecone
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import (
    OPENAI_API_KEY,
    PINECONE_API_KEY,
    INDEX_NAME
)

# =====================================
# INIT
# =====================================

client = OpenAI(api_key=OPENAI_API_KEY)

pc = Pinecone(api_key=PINECONE_API_KEY)
index = pc.Index(INDEX_NAME)

DATA_DIR = "data"


# =====================================
# EMBEDDING
# =====================================

def get_embedding(text: str):
    response = client.embeddings.create(
        model="text-embedding-3-small",
        input=text
    )
    return response.data[0].embedding


# =====================================
# PPT EXTRACTION
# =====================================

def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)

    slides_text = []

    for i, slide in enumerate(prs.slides):
        content = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content.append(shape.text.strip())

        if content:
            slides_text.append({
                "text": f"[اسلاید {i+1}]\n" + "\n".join(content),
                "slide": i + 1
            })

    return slides_text


# =====================================
# BUILD VECTOR DB
# =====================================

def build_vector_db():

    all_docs = []

    for filename in os.listdir(DATA_DIR):

        if not filename.endswith(".pptx"):
            continue

        filepath = os.path.join(DATA_DIR, filename)
        print(f"Processing: {filename}")

        slides = extract_text_from_pptx(filepath)

        for s in slides:
            all_docs.append({
                "text": s["text"],
                "slide": s["slide"],
                "source": filename
            })

    if not all_docs:
        print("No PPTX files found")
        return

    # -----------------------------
    # CHUNKING
    # -----------------------------
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100
    )

    chunks = splitter.create_documents(
        [d["text"] for d in all_docs]
    )

    print(f"Chunks created: {len(chunks)}")

    vectors = []

    # -----------------------------
    # VECTOR BUILD
    # -----------------------------
    for i, chunk in enumerate(chunks):

        print(f"Embedding {i+1}/{len(chunks)}")

        embedding = get_embedding(chunk.page_content)

        vectors.append({
            "id": f"ppt_{i}",
            "values": embedding,
            "metadata": {
                "text": chunk.page_content,
                "source": "tank_storage",
                "chunk_id": i
            }
        })

        # batch upload
        if len(vectors) >= 50:
            index.upsert(vectors=vectors)
            vectors = []

    if vectors:
        index.upsert(vectors=vectors)

    print("Upload completed successfully")


# =====================================
# MAIN
# =====================================

if __name__ == "__main__":
    build_vector_db()